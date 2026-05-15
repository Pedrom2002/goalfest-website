"""
Generates Fanzone Lisboa 2026 sponsor pitch deck (.pptx)
Colors match brand: #0d0d0d bg, #FFD700 gold, #C8102E red-pt
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import os

# Brand colors
C_BG = RGBColor(0x0D, 0x0D, 0x0D)
C_SURFACE = RGBColor(0x1A, 0x1A, 0x1A)
C_GOLD = RGBColor(0xFF, 0xD7, 0x00)
C_RED = RGBColor(0xC8, 0x10, 0x2E)
C_WHITE = RGBColor(0xF5, 0xF5, 0xF5)
C_MUTED = RGBColor(0x9C, 0xA3, 0xAF)
C_GREEN = RGBColor(0x00, 0x66, 0x00)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


def new_prs():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_slide(prs):
    blank_layout = prs.slide_layouts[6]  # completely blank
    return prs.slides.add_slide(blank_layout)


def fill_bg(slide, color=C_BG):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_text(slide, text, left, top, width, height,
             font_size=24, bold=False, color=C_WHITE,
             align=PP_ALIGN.LEFT, italic=False):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


def add_paragraph(tf, text, font_size=16, bold=False,
                  color=C_WHITE, align=PP_ALIGN.LEFT, space_before=6):
    from pptx.util import Pt as pt
    from pptx.oxml.ns import qn
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return p


def gold_bar(slide, top=Inches(0.12)):
    add_rect(slide, 0, top, SLIDE_W, Pt(4), C_GOLD)


def slide_number_tag(slide, num, total=14):
    add_text(slide, f"{num}/{total}",
             SLIDE_W - Inches(1.2), SLIDE_H - Inches(0.45),
             Inches(1), Inches(0.35),
             font_size=11, color=C_MUTED, align=PP_ALIGN.RIGHT)


# ─── SLIDE 1: CAPA ───────────────────────────────────────────────────────────
def slide_01_capa(prs):
    sl = blank_slide(prs)
    fill_bg(sl)

    # Gradient-like accent strips
    add_rect(sl, 0, 0, Inches(0.6), SLIDE_H, C_RED)
    add_rect(sl, SLIDE_W - Inches(0.6), 0, Inches(0.6), SLIDE_H, C_GREEN)
    gold_bar(sl, top=Inches(0.22))
    gold_bar(sl, top=SLIDE_H - Inches(0.3))

    add_text(sl, "FANZONE LISBOA 2026",
             Inches(1), Inches(1.8), Inches(11), Inches(1.4),
             font_size=52, bold=True, color=C_GOLD, align=PP_ALIGN.CENTER)

    add_text(sl, "O Maior Evento Público do Mundial em Portugal",
             Inches(1), Inches(3.2), Inches(11), Inches(0.9),
             font_size=24, color=C_WHITE, align=PP_ALIGN.CENTER)

    add_text(sl, "Proposta de Parceria · Junho–Julho 2026",
             Inches(1), Inches(4.2), Inches(11), Inches(0.7),
             font_size=16, color=C_MUTED, align=PP_ALIGN.CENTER, italic=True)

    add_text(sl, "fanzonelisboa.pt",
             Inches(1), Inches(6.4), Inches(11), Inches(0.5),
             font_size=14, color=C_GOLD, align=PP_ALIGN.CENTER)


# ─── SLIDE 2: A OPORTUNIDADE ─────────────────────────────────────────────────
def slide_02_oportunidade(prs):
    sl = blank_slide(prs)
    fill_bg(sl)
    gold_bar(sl)
    slide_number_tag(sl, 2)

    add_text(sl, "A OPORTUNIDADE",
             Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
             font_size=11, color=C_GOLD, bold=True)

    add_text(sl, "O Mundial de 2026 é o maior da história",
             Inches(0.8), Inches(1.1), Inches(11), Inches(1.0),
             font_size=34, bold=True, color=C_WHITE)

    bullets = [
        ("🏆", "48 seleções  |  104 jogos  |  3 países anfitriões (EUA, Canadá, México)"),
        ("🇵🇹", "Portugal qualificado — jogos transmitidos em direto"),
        ("📺", "Audiência global estimada: 5 mil milhões de espectadores"),
        ("🏙️", "Lisboa posiciona-se como capital do futebol europeu neste verão"),
        ("🧑‍🤝‍🧑", "Geração de comunidade, emoção partilhada e consumo ao vivo"),
    ]

    y = Inches(2.2)
    for icon, text in bullets:
        add_rect(sl, Inches(0.8), y + Inches(0.08), Inches(0.05), Inches(0.3), C_GOLD)
        add_text(sl, f"{icon}  {text}",
                 Inches(1.0), y, Inches(11), Inches(0.5),
                 font_size=16, color=C_WHITE)
        y += Inches(0.62)


# ─── SLIDE 3: O EVENTO ───────────────────────────────────────────────────────
def slide_03_evento(prs):
    sl = blank_slide(prs)
    fill_bg(sl)
    gold_bar(sl)
    slide_number_tag(sl, 3)

    add_text(sl, "O EVENTO",
             Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
             font_size=11, color=C_GOLD, bold=True)

    add_text(sl, "Fanzone Lisboa 2026",
             Inches(0.8), Inches(1.0), Inches(10), Inches(0.9),
             font_size=38, bold=True, color=C_WHITE)

    # Info cards (3 side by side)
    cards = [
        ("📍", "LOCALIZAÇÃO", "Lisboa, Portugal\n(centro da cidade)"),
        ("📅", "DATAS", "Junho – Julho 2026\nDurante todo o torneio"),
        ("👥", "CAPACIDADE", "500–1.000 pessoas\npor jogo"),
    ]

    cx = Inches(0.8)
    for icon, title, body in cards:
        add_rect(sl, cx, Inches(2.1), Inches(3.6), Inches(2.8), C_SURFACE)
        add_text(sl, icon,
                 cx + Inches(0.2), Inches(2.25), Inches(0.8), Inches(0.6),
                 font_size=28)
        add_text(sl, title,
                 cx + Inches(0.2), Inches(2.9), Inches(3.2), Inches(0.45),
                 font_size=13, bold=True, color=C_GOLD)
        add_text(sl, body,
                 cx + Inches(0.2), Inches(3.35), Inches(3.2), Inches(1.2),
                 font_size=15, color=C_WHITE)
        cx += Inches(4.1)

    add_text(sl, "Transmissão em ecrãs gigantes · Ambiente de festival · Entrada gratuita",
             Inches(0.8), Inches(5.3), Inches(11.5), Inches(0.6),
             font_size=15, color=C_MUTED, align=PP_ALIGN.CENTER, italic=True)


# ─── SLIDE 4: A AUDIÊNCIA ────────────────────────────────────────────────────
def slide_04_audiencia(prs):
    sl = blank_slide(prs)
    fill_bg(sl)
    gold_bar(sl)
    slide_number_tag(sl, 4)

    add_text(sl, "A AUDIÊNCIA",
             Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
             font_size=11, color=C_GOLD, bold=True)

    add_text(sl, "Quem vai estar na Fanzone",
             Inches(0.8), Inches(1.0), Inches(10), Inches(0.9),
             font_size=36, bold=True, color=C_WHITE)

    # Big stats
    stats = [
        ("15.000–20.000", "visitantes únicos\nest. ao longo do torneio"),
        ("70%", "idades 18–45 anos\npúblico jovem e ativo"),
        ("40%", "turistas estrangeiros\nem Lisboa no período"),
        ("60%+", "presença nas redes sociais\npúblico digital-first"),
    ]

    cx = Inches(0.5)
    for num, label in stats:
        add_rect(sl, cx, Inches(2.1), Inches(2.9), Inches(2.5), C_SURFACE)
        add_text(sl, num,
                 cx + Inches(0.15), Inches(2.25), Inches(2.6), Inches(1.0),
                 font_size=32, bold=True, color=C_GOLD, align=PP_ALIGN.CENTER)
        add_text(sl, label,
                 cx + Inches(0.15), Inches(3.25), Inches(2.6), Inches(0.9),
                 font_size=13, color=C_MUTED, align=PP_ALIGN.CENTER)
        cx += Inches(3.15)

    add_text(sl, "Perfil: adeptos de futebol, famílias, grupos de amigos, turistas — consumidores ativos",
             Inches(0.8), Inches(5.1), Inches(11.5), Inches(0.6),
             font_size=14, color=C_WHITE, align=PP_ALIGN.CENTER, italic=True)


# ─── SLIDE 5: WEBSITE & DIGITAL ──────────────────────────────────────────────
def slide_05_website(prs):
    sl = blank_slide(prs)
    fill_bg(sl)
    gold_bar(sl)
    slide_number_tag(sl, 5)

    add_text(sl, "PRESENÇA DIGITAL",
             Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
             font_size=11, color=C_GOLD, bold=True)

    add_text(sl, "fanzonelisboa.pt — o hub digital do evento",
             Inches(0.8), Inches(1.0), Inches(11), Inches(0.9),
             font_size=32, bold=True, color=C_WHITE)

    features = [
        ("🌐", "Website oficial", "Calendário de jogos, reservas, informações — disponível em PT e EN"),
        ("📱", "Mobile-first", "Optimizado para smartphone: a maioria do tráfego vem de mobile"),
        ("📊", "Métricas estimadas", "5.000+ visitas/mês em pico · 500+ reservas esperadas por jogo"),
        ("🔗", "SEO & redes sociais", "Presença orgânica no Google + perfis activos em Instagram/TikTok"),
        ("🏷️", "Secção de Patrocinadores", "Logo e link da sua marca visível em todas as páginas do site"),
    ]

    y = Inches(2.1)
    for icon, title, desc in features:
        add_rect(sl, Inches(0.8), y + Inches(0.05), Inches(0.05), Inches(0.38), C_GOLD)
        add_text(sl, f"{icon}  {title}",
                 Inches(1.1), y, Inches(4.5), Inches(0.48),
                 font_size=15, bold=True, color=C_GOLD)
        add_text(sl, desc,
                 Inches(5.8), y, Inches(7), Inches(0.48),
                 font_size=14, color=C_WHITE)
        y += Inches(0.9)


# ─── SLIDE 6: PORQUE LISBOA ──────────────────────────────────────────────────
def slide_06_lisboa(prs):
    sl = blank_slide(prs)
    fill_bg(sl)
    gold_bar(sl)
    slide_number_tag(sl, 6)

    add_text(sl, "PORQUE LISBOA",
             Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
             font_size=11, color=C_GOLD, bold=True)

    add_text(sl, "A cidade certa, no momento certo",
             Inches(0.8), Inches(1.0), Inches(11), Inches(0.9),
             font_size=36, bold=True, color=C_WHITE)

    points = [
        ("🏆", "Capital europeia do futebol",
         "Com Portugal no Mundial, Lisboa será o epicentro da celebração europeia durante 6 semanas"),
        ("✈️", "Turismo recorde",
         "Lisboa recebe 10M+ turistas por ano — o verão de 2026 promete números históricos"),
        ("📡", "Cobertura mediática",
         "Media nacional e internacional em Lisboa durante o torneio"),
        ("💰", "Poder de compra",
         "Perfil de consumidor jovem, urbano e com hábitos de consumo elevados"),
    ]

    y = Inches(2.1)
    for icon, title, desc in points:
        add_rect(sl, Inches(0.8), y, Inches(11.5), Inches(1.0), C_SURFACE)
        add_text(sl, f"{icon}  {title}",
                 Inches(1.0), y + Inches(0.08), Inches(4.5), Inches(0.45),
                 font_size=15, bold=True, color=C_GOLD)
        add_text(sl, desc,
                 Inches(5.5), y + Inches(0.08), Inches(6.5), Inches(0.8),
                 font_size=14, color=C_WHITE)
        y += Inches(1.15)


# ─── SLIDE 7: VISIBILIDADE DA MARCA ──────────────────────────────────────────
def slide_07_visibilidade(prs):
    sl = blank_slide(prs)
    fill_bg(sl)
    gold_bar(sl)
    slide_number_tag(sl, 7)

    add_text(sl, "A VOSSA MARCA NA FANZONE",
             Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
             font_size=11, color=C_GOLD, bold=True)

    add_text(sl, "Visibilidade onde importa",
             Inches(0.8), Inches(1.0), Inches(11), Inches(0.9),
             font_size=36, bold=True, color=C_WHITE)

    touchpoints = [
        ("🖥️", "Website", "Logo + link na secção de patrocinadores · visível em todas as páginas"),
        ("📸", "Redes Sociais", "Menções em posts e stories · conteúdo co-branded em jogos Portugal"),
        ("📣", "Comunicações", "Logo em newsletters e comunicados de imprensa do evento"),
        ("🎯", "Público qualificado", "Alcance direto a consumidores no contexto emocional mais poderoso"),
    ]

    y = Inches(2.1)
    for icon, title, desc in touchpoints:
        add_rect(sl, Inches(0.8), y + Inches(0.05), Inches(0.05), Inches(0.38), C_RED)
        add_text(sl, f"{icon}  {title}",
                 Inches(1.1), y, Inches(4.0), Inches(0.48),
                 font_size=15, bold=True, color=C_WHITE)
        add_text(sl, desc,
                 Inches(5.5), y, Inches(7), Inches(0.48),
                 font_size=14, color=C_MUTED)
        y += Inches(0.95)

    add_rect(sl, Inches(0.8), Inches(6.1), Inches(11.5), Inches(0.9), C_RED)
    add_text(sl, "💡  O patrocinador associa-se à emoção do Mundial — o evento mais visto do planeta",
             Inches(1.0), Inches(6.15), Inches(11), Inches(0.7),
             font_size=15, bold=True, color=C_WHITE)


# ─── SLIDE 8: REDES SOCIAIS ───────────────────────────────────────────────────
def slide_08_social(prs):
    sl = blank_slide(prs)
    fill_bg(sl)
    gold_bar(sl)
    slide_number_tag(sl, 8)

    add_text(sl, "VISIBILIDADE DIGITAL",
             Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
             font_size=11, color=C_GOLD, bold=True)

    add_text(sl, "Alcance orgânico estimado: 50.000+ pessoas",
             Inches(0.8), Inches(1.0), Inches(11), Inches(0.9),
             font_size=32, bold=True, color=C_WHITE)

    channels = [
        ("📸 Instagram", "Posts de jogos, stories ao vivo, reels do evento\nMenção à marca nos posts patrocinados"),
        ("🎵 TikTok", "Conteúdo viral de momentos do evento\nTags e menções da marca parceira"),
        ("👥 Facebook", "Eventos públicos, transmissões ao vivo\nAlcance geográfico Lisboa + região"),
    ]

    cx = Inches(0.8)
    for title, desc in channels:
        add_rect(sl, cx, Inches(2.2), Inches(3.7), Inches(2.8), C_SURFACE)
        add_text(sl, title,
                 cx + Inches(0.2), Inches(2.4), Inches(3.3), Inches(0.55),
                 font_size=18, bold=True, color=C_GOLD)
        add_text(sl, desc,
                 cx + Inches(0.2), Inches(3.0), Inches(3.3), Inches(1.6),
                 font_size=14, color=C_WHITE)
        cx += Inches(4.1)

    add_text(sl, "Estratégia: conteúdo antes, durante e após cada jogo · pico de engagement nos jogos de Portugal",
             Inches(0.8), Inches(5.4), Inches(11.5), Inches(0.7),
             font_size=14, color=C_MUTED, align=PP_ALIGN.CENTER, italic=True)


# ─── SLIDE 9: PACOTES ─────────────────────────────────────────────────────────
def slide_09_pacotes(prs):
    sl = blank_slide(prs)
    fill_bg(sl)
    gold_bar(sl)
    slide_number_tag(sl, 9)

    add_text(sl, "PACOTES DE PATROCÍNIO",
             Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
             font_size=11, color=C_GOLD, bold=True)

    add_text(sl, "Escolha o nível de visibilidade",
             Inches(0.8), Inches(1.0), Inches(11), Inches(0.8),
             font_size=34, bold=True, color=C_WHITE)

    tiers = [
        (C_GOLD,    "🥇 OURO",    "Patrocinador Principal",  "€5.000"),
        (C_MUTED,   "🥈 PRATA",   "Parceiro Oficial",        "€2.500"),
        (RGBColor(0xCD, 0x7F, 0x32), "🥉 BRONZE", "Parceiro de Apoio", "€1.000"),
    ]

    cx = Inches(0.6)
    for color, tier, subtitle, price in tiers:
        add_rect(sl, cx, Inches(2.0), Inches(3.8), Inches(3.8), C_SURFACE)
        add_rect(sl, cx, Inches(2.0), Inches(3.8), Inches(0.08), color)
        add_text(sl, tier,
                 cx + Inches(0.2), Inches(2.15), Inches(3.4), Inches(0.65),
                 font_size=20, bold=True, color=color)
        add_text(sl, subtitle,
                 cx + Inches(0.2), Inches(2.8), Inches(3.4), Inches(0.5),
                 font_size=14, color=C_MUTED)
        add_text(sl, price,
                 cx + Inches(0.2), Inches(3.35), Inches(3.4), Inches(0.8),
                 font_size=32, bold=True, color=C_WHITE)
        add_text(sl, "ver detalhes no slide seguinte →",
                 cx + Inches(0.2), Inches(4.3), Inches(3.4), Inches(0.45),
                 font_size=12, color=color, italic=True)
        cx += Inches(4.2)

    add_text(sl, "* Valores indicativos · Pacotes personalizados disponíveis sob consulta",
             Inches(0.8), Inches(6.3), Inches(11.5), Inches(0.5),
             font_size=12, color=C_MUTED, align=PP_ALIGN.CENTER, italic=True)


# ─── SLIDE 10: TABELA COMPARATIVA ─────────────────────────────────────────────
def slide_10_tabela(prs):
    sl = blank_slide(prs)
    fill_bg(sl)
    gold_bar(sl)
    slide_number_tag(sl, 10)

    add_text(sl, "O QUE ESTÁ INCLUÍDO",
             Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
             font_size=11, color=C_GOLD, bold=True)

    add_text(sl, "Comparação de benefícios por pacote",
             Inches(0.8), Inches(1.0), Inches(11), Inches(0.8),
             font_size=30, bold=True, color=C_WHITE)

    rows = [
        ("Benefício",                             "🥇 Ouro",   "🥈 Prata",  "🥉 Bronze"),
        ("Logo no website (todas as páginas)",    "✅ Grande",  "✅ Médio",  "✅ Pequeno"),
        ("Menções em redes sociais",              "✅ 10+",     "✅ 5",      "✅ 2"),
        ("Post dedicado à marca",                 "✅ Sim",     "❌ Não",    "❌ Não"),
        ("Logo em newsletters do evento",         "✅ Sim",     "✅ Sim",    "❌ Não"),
        ("Título de patrocinador principal",      "✅ Sim",     "❌ Não",    "❌ Não"),
        ("Conteúdo co-branded nos jogos PT",      "✅ Todos",   "✅ 2 jogos","❌ Não"),
        ("Menção em comunicados de imprensa",     "✅ Sim",     "✅ Sim",    "❌ Não"),
    ]

    col_widths = [Inches(5.5), Inches(2.2), Inches(2.2), Inches(2.2)]
    col_x = [Inches(0.5), Inches(6.2), Inches(8.5), Inches(10.8)]
    row_h = Inches(0.52)
    start_y = Inches(2.0)

    for r_idx, row in enumerate(rows):
        bg = C_SURFACE if r_idx % 2 == 0 else C_BG
        if r_idx == 0:
            bg = RGBColor(0x11, 0x11, 0x11)
        add_rect(sl, Inches(0.5), start_y + r_idx * row_h, Inches(12.3), row_h, bg)
        for c_idx, cell in enumerate(row):
            is_header = r_idx == 0
            col_color = C_GOLD if is_header else (C_WHITE if c_idx == 0 else C_MUTED)
            add_text(sl, cell,
                     col_x[c_idx], start_y + r_idx * row_h + Inches(0.08),
                     col_widths[c_idx], row_h - Inches(0.1),
                     font_size=13, bold=is_header, color=col_color)


# ─── SLIDE 11: CASOS DE SUCESSO ───────────────────────────────────────────────
def slide_11_casos(prs):
    sl = blank_slide(prs)
    fill_bg(sl)
    gold_bar(sl)
    slide_number_tag(sl, 11)

    add_text(sl, "REFERÊNCIAS DE MERCADO",
             Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
             font_size=11, color=C_GOLD, bold=True)

    add_text(sl, "Fanzones que funcionam — e as marcas que ganham",
             Inches(0.8), Inches(1.0), Inches(11), Inches(0.9),
             font_size=30, bold=True, color=C_WHITE)

    cases = [
        ("UEFA Euro 2024", "Berlim, Alemanha",
         "Fanzone oficial com 70.000 pessoas/dia\nPatrocinadores: Adidas, Coca-Cola, Lidl\nROI: cobertura em 180 países"),
        ("Fanzone Lisboa Euro 2016", "Lisboa, PT",
         "15.000 adeptos em noites de jogos de Portugal\nMarcas locais duplicaram vendas no período\nViral nas redes sociais portuguesas"),
        ("FIFA World Cup 2022", "Zonas de fan em todo o mundo",
         "Marcas de bebidas reportaram +40% de vendas\nAssociação à emoção do Mundial = recall de marca elevado"),
    ]

    cx = Inches(0.5)
    for event, location, desc in cases:
        add_rect(sl, cx, Inches(2.1), Inches(3.9), Inches(3.5), C_SURFACE)
        add_rect(sl, cx, Inches(2.1), Inches(3.9), Inches(0.06), C_GOLD)
        add_text(sl, event,
                 cx + Inches(0.2), Inches(2.2), Inches(3.5), Inches(0.55),
                 font_size=16, bold=True, color=C_GOLD)
        add_text(sl, location,
                 cx + Inches(0.2), Inches(2.75), Inches(3.5), Inches(0.4),
                 font_size=13, color=C_MUTED, italic=True)
        add_text(sl, desc,
                 cx + Inches(0.2), Inches(3.2), Inches(3.5), Inches(2.0),
                 font_size=13, color=C_WHITE)
        cx += Inches(4.2)


# ─── SLIDE 12: A EQUIPA ───────────────────────────────────────────────────────
def slide_12_equipa(prs):
    sl = blank_slide(prs)
    fill_bg(sl)
    gold_bar(sl)
    slide_number_tag(sl, 12)

    add_text(sl, "QUEM SOMOS",
             Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
             font_size=11, color=C_GOLD, bold=True)

    add_text(sl, "A equipa por trás da Fanzone Lisboa",
             Inches(0.8), Inches(1.0), Inches(11), Inches(0.9),
             font_size=34, bold=True, color=C_WHITE)

    add_text(sl, "Somos um grupo de apaixonados por futebol e eventos em Lisboa, com experiência em organização de eventos desportivos públicos e presença digital.",
             Inches(0.8), Inches(2.0), Inches(11.5), Inches(1.0),
             font_size=16, color=C_MUTED)

    members = [
        ("Organização", "Experiência em eventos\ndesportivos públicos\nem Lisboa"),
        ("Digital & Marketing", "Website, redes sociais,\nstratégia de conteúdo\ne comunicação"),
        ("Operações & Local", "Logística, equipamento\naudiovisual, segurança\ne catering"),
    ]

    cx = Inches(0.8)
    for role, desc in members:
        add_rect(sl, cx, Inches(3.2), Inches(3.6), Inches(2.5), C_SURFACE)
        add_rect(sl, cx, Inches(3.2), Inches(3.6), Inches(0.06), C_RED)
        add_text(sl, role,
                 cx + Inches(0.2), Inches(3.35), Inches(3.2), Inches(0.5),
                 font_size=15, bold=True, color=C_WHITE)
        add_text(sl, desc,
                 cx + Inches(0.2), Inches(3.9), Inches(3.2), Inches(1.4),
                 font_size=13, color=C_MUTED)
        cx += Inches(4.1)

    add_text(sl, "Abertos a apresentar a equipa pessoalmente em reunião — disponíveis em Lisboa",
             Inches(0.8), Inches(6.1), Inches(11.5), Inches(0.6),
             font_size=14, color=C_GOLD, align=PP_ALIGN.CENTER, italic=True)


# ─── SLIDE 13: PRÓXIMOS PASSOS ────────────────────────────────────────────────
def slide_13_timeline(prs):
    sl = blank_slide(prs)
    fill_bg(sl)
    gold_bar(sl)
    slide_number_tag(sl, 13)

    add_text(sl, "PRÓXIMOS PASSOS",
             Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
             font_size=11, color=C_GOLD, bold=True)

    add_text(sl, "Timeline de parceria",
             Inches(0.8), Inches(1.0), Inches(11), Inches(0.9),
             font_size=36, bold=True, color=C_WHITE)

    steps = [
        ("Maio 2026",    C_RED,   "Fecho de parcerias",
         "Prazo para confirmar patrocínios e integrar logos em todos os materiais"),
        ("Junho 2026",   C_GOLD,  "Início da visibilidade digital",
         "Website live, redes sociais activas, primeira campanha de comunicação"),
        ("11 Jun 2026",  C_GREEN, "Arranque do Mundial",
         "Primeiro jogo transmitido na Fanzone — cobertura mediática máxima"),
        ("Jun–Jul 2026", C_WHITE, "Período de patrocínio activo",
         "Visibilidade contínua durante todos os jogos transmitidos"),
        ("Após torneio", C_MUTED, "Relatório de impacto",
         "Métricas de alcance, fotos e relatório enviado a todos os patrocinadores"),
    ]

    y = Inches(2.05)
    for date, color, title, desc in steps:
        add_rect(sl, Inches(0.8), y + Inches(0.1), Inches(0.08), Inches(0.55), color)
        add_text(sl, date,
                 Inches(1.1), y, Inches(1.8), Inches(0.45),
                 font_size=13, bold=True, color=color)
        add_text(sl, title,
                 Inches(3.0), y, Inches(3.8), Inches(0.45),
                 font_size=14, bold=True, color=C_WHITE)
        add_text(sl, desc,
                 Inches(7.0), y, Inches(6.0), Inches(0.55),
                 font_size=13, color=C_MUTED)
        y += Inches(0.85)

    add_rect(sl, Inches(0.8), Inches(6.55), Inches(11.5), Inches(0.6), C_RED)
    add_text(sl, "⚠️  Vagas limitadas — os patrocinadores Ouro são exclusivos (máx. 1 por categoria de produto)",
             Inches(1.0), Inches(6.6), Inches(11.0), Inches(0.5),
             font_size=13, bold=True, color=C_WHITE)


# ─── SLIDE 14: CTA / CONTACTO ─────────────────────────────────────────────────
def slide_14_cta(prs):
    sl = blank_slide(prs)
    fill_bg(sl)
    add_rect(sl, 0, 0, Inches(0.6), SLIDE_H, C_RED)
    add_rect(sl, SLIDE_W - Inches(0.6), 0, Inches(0.6), SLIDE_H, C_GREEN)
    gold_bar(sl, top=Inches(0.22))
    gold_bar(sl, top=SLIDE_H - Inches(0.3))

    add_text(sl, "VAMOS CONSTRUIR ISTO JUNTOS",
             Inches(1), Inches(1.2), Inches(11), Inches(1.2),
             font_size=40, bold=True, color=C_GOLD, align=PP_ALIGN.CENTER)

    add_text(sl, "Junte a sua marca à maior celebração de futebol em Lisboa",
             Inches(1), Inches(2.5), Inches(11), Inches(0.7),
             font_size=20, color=C_WHITE, align=PP_ALIGN.CENTER)

    contact_info = [
        "📧  geral@fanzonelisboa.pt",
        "🌐  fanzonelisboa.pt",
        "📱  Instagram: @fanzonelisboa",
    ]

    y = Inches(3.5)
    for info in contact_info:
        add_text(sl, info,
                 Inches(2), y, Inches(9), Inches(0.55),
                 font_size=18, color=C_WHITE, align=PP_ALIGN.CENTER)
        y += Inches(0.65)

    add_rect(sl, Inches(3.5), Inches(5.6), Inches(6), Inches(0.8), C_RED)
    add_text(sl, "MARCAR REUNIÃO",
             Inches(3.5), Inches(5.65), Inches(6), Inches(0.7),
             font_size=20, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    add_text(sl, "Disponíveis para reunião presencial em Lisboa ou videochamada",
             Inches(1), Inches(6.65), Inches(11), Inches(0.55),
             font_size=13, color=C_MUTED, align=PP_ALIGN.CENTER, italic=True)


def main():
    prs = new_prs()

    slide_01_capa(prs)
    slide_02_oportunidade(prs)
    slide_03_evento(prs)
    slide_04_audiencia(prs)
    slide_05_website(prs)
    slide_06_lisboa(prs)
    slide_07_visibilidade(prs)
    slide_08_social(prs)
    slide_09_pacotes(prs)
    slide_10_tabela(prs)
    slide_11_casos(prs)
    slide_12_equipa(prs)
    slide_13_timeline(prs)
    slide_14_cta(prs)

    out_dir = os.path.join(os.path.dirname(__file__), "..", "docs")
    os.makedirs(out_dir, exist_ok=True)
    out_path = os.path.join(out_dir, "Fanzone_Lisboa_2026_Pitch_Deck.pptx")
    prs.save(out_path)
    print(f"Saved: {os.path.abspath(out_path)}")


if __name__ == "__main__":
    main()
